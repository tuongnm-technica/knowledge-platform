import io
import structlog

log = structlog.get_logger()


class FileParser:

    def parse(self, filename: str, data: bytes) -> str:
        ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        try:
            if ext == ".pdf":
                return self._parse_pdf(data)
            elif ext == ".docx":
                return self._parse_docx(data)
            elif ext == ".xlsx":
                return self._parse_xlsx(data)
            elif ext == ".pptx":
                return self._parse_pptx(data)
            elif ext in [".txt", ".md"]:
                return data.decode("utf-8", errors="ignore")
            else:
                return ""
        except Exception as e:
            log.error("file_parser.failed", filename=filename, error=str(e))
            return ""

    def _parse_pdf(self, data: bytes) -> str:
        import fitz  # pymupdf
        doc  = fitz.open(stream=data, filetype="pdf")
        text = "\n\n".join(page.get_text() for page in doc)
        doc.close()
        return text.strip()

    def _parse_docx(self, data: bytes) -> str:
        from docx import Document
        doc   = Document(io.BytesIO(data))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        # Đọc cả tables
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n\n".join(parts)

    def _parse_xlsx(self, data: bytes) -> str:
        import openpyxl
        wb    = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            parts.append(f"Sheet: {sheet}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts)

    def _parse_pptx(self, data: bytes) -> str:
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"Slide {i}:\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)